from io import BytesIO

import pdfplumber
import pandas as pd
from tempfile import NamedTemporaryFile
import os
import docx
from openpyxl.reader.excel import load_workbook
from openpyxl.workbook import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter, landscape, A4, portrait
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from docx import Document
from docx.oxml import OxmlElement


def pdf_to_excel(pdf_path):
    excel_path = os.path.splitext(pdf_path)[0] + '.xlsx'
    
    with pdfplumber.open(pdf_path) as pdf:
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            for i, page in enumerate(pdf.pages):
                table = page.extract_table()
                if table:
                    df = pd.DataFrame(table[1:], columns=table[0])
                    sheet_name = f'Page_{i + 1}'
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
    print(f"PDF Conversion completed: {excel_path}")

    return excel_path

def docx_to_excel(docx_path):
    excel_path = os.path.splitext(docx_path)[0] + '.xlsx'
    
    doc = docx.Document(docx_path)
    data = []

    for table in doc.tables:
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                # Add cell text to row data
                row_data.append(cell.text)
            data.append(row_data)

    if data:
        # Adjust for merged cells (if necessary)
        max_cols = max(len(row) for row in data)
        for row in data:
            while len(row) < max_cols:
                row.append('')

        df = pd.DataFrame(data[1:], columns=data[0])
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Sheet1', index=False)
    print(f"DOCX Conversion completed: {excel_path}")
    return excel_path

def ppt_to_excel(ppt_path):
    excel_path = os.path.splitext(ppt_path)[0] + '.xlsx'

    # Load the PowerPoint presentation
    prs = Presentation(ppt_path)

    # Create a new Excel workbook and select the active worksheet
    wb = Workbook()
    ws = wb.active
    ws.title = "Tables"

    # Iterate over slides in the presentation
    for slide_index, slide in enumerate(prs.slides):
        # Iterate over shapes in the slide
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                # Get the table data and write it to the Excel sheet
                for row_index, row in enumerate(table.rows):
                    for col_index, cell in enumerate(row.cells):
                        cell_value = cell.text.strip()
                        ws.cell(row=row_index + 1, column=col_index + 1, value=cell_value)
                # Move to a new sheet for the next table
                ws = wb.create_sheet(title=f"Table_Slide_{slide_index + 1}")

    # Remove the default empty sheet if there is more than one sheet
    if len(wb.sheetnames) > 1:
        wb.remove(wb['Tables'])

    wb.save(excel_path)
    print(f"PPT Conversion completed: {excel_path}")

    return excel_path


def excel_to_pdf(excel_path, pdf_path):
    workbook = load_workbook(excel_path)
    sheet_names = workbook.sheetnames
    doc = SimpleDocTemplate(pdf_path, pagesize=letter)
    elements = []
    styles = getSampleStyleSheet()
    title_style = styles['Title']
    body_style = styles['BodyText']

    for sheet_name in sheet_names:
        sheet = workbook[sheet_name]
        df = pd.DataFrame(sheet.values)

        if not df.empty:
            # Check for merged cells and adjust the DataFrame accordingly
            for merge in sheet.merged_cells.ranges:
                merged_cell = merge.coord
                start_row, start_col, end_row, end_col = merge.min_row - 1, merge.min_col - 1, merge.max_row - 1, merge.max_col - 1
                merge_value = df.iloc[start_row, start_col]
                df.iloc[start_row:end_row + 1, start_col:end_col + 1] = None
                df.iloc[start_row, start_col] = merge_value

            elements.append(Paragraph(sheet_name, title_style))
            elements.append(Spacer(1, 12))  # Add space after the title

            data = [df.columns.tolist()] + df.fillna('').values.tolist()
            if data and len(data[0]) > 0:
                table = Table(data)

                # Setting table styles for better visual representation
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('ALIGN', (0, 1), (-1, -1), 'CENTER'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ]))

                elements.append(table)
                elements.append(Spacer(1, 24))  # Add space after the table

    if elements:
        doc.build(elements)
        print(f"PDF created: {pdf_path}")
    else:
        print("No valid data found in the Excel file to create a PDF.")

    return pdf_path

def set_cell_border(cell, **kwargs):
    tc = cell._element
    for key, value in kwargs.items():
        tc.get_or_add_tcPr().append(OxmlElement(key))
        element = tc.tcPr[-1]
        element.append(OxmlElement(value))


def excel_to_docx(excel_path, docx_path):
    """
    Convert an Excel workbook to a DOCX document with tables.
    """
    # Load Excel workbook
    workbook = load_workbook(excel_path, data_only=True)
    sheet_names = workbook.sheetnames

    # Create a new Document object
    doc = Document()

    for sheet_name in sheet_names:
        # Add sheet name as heading
        doc.add_heading(sheet_name, level=1)

        # Create a table for the current sheet
        sheet = workbook[sheet_name]
        table = doc.add_table(rows=sheet.max_row, cols=sheet.max_column)

        # Set table style
        table.style = 'Table Grid'

        # Iterate through each cell in the sheet
        for row in sheet.iter_rows():
            for cell in row:
                table_cell = table.cell(cell.row - 1, cell.column - 1)
                value = str(cell.value) if cell.value is not None else ''
                table_cell.text = value

    # Save the document
    doc.save(docx_path)
    print(f"DOCX created: {docx_path}")

    return docx_path


def docx_to_pdf(docx_path, pdf_path, orientation):
    # Load DOCX file
    doc = Document(docx_path)

    # Determine page size based on orientation
    if orientation == 'landscape':
        pagesize = landscape(A4)
    else:
        pagesize = portrait(A4)

    # Create a BytesIO buffer for the PDF
    buffer = BytesIO()

    # Create a canvas with specified page size
    c = canvas.Canvas(buffer, pagesize=pagesize)

    # Set margins
    left_margin = 50
    right_margin = 50
    top_margin = 50
    bottom_margin = 50

    # Calculate usable width and height considering margins
    usable_width = pagesize[0] - left_margin - right_margin
    usable_height = pagesize[1] - top_margin - bottom_margin

    # Styles
    styles = getSampleStyleSheet()
    normal_style = styles["Normal"]
    heading_style = styles["Heading1"]

    try:
        # Iterate through paragraphs and tables in DOCX
        y_position = pagesize[1] - top_margin
        for para in doc.paragraphs:
            # Draw each paragraph
            ptext = para.text
            p = Paragraph(ptext, normal_style)
            p.wrapOn(c, usable_width, 20)

            # Check if adding this paragraph will exceed the page height
            if y_position - p.height < bottom_margin:
                c.showPage()  # Start a new page
                c.setFont("Helvetica", 12)
                y_position = pagesize[1] - top_margin

            p.drawOn(c, left_margin, y_position - p.height)

            # Adjust Y position
            y_position -= p.height + 10  # Adding some space between paragraphs

        # Draw tables
        for table in doc.tables:
            # Assuming 1st row is header and rest are data
            data = []
            for i, row in enumerate(table.rows):
                if i == 0:
                    continue  # skip header row
                data.append([cell.text for cell in row.cells])

            # Convert table data to PDF table
            pdf_table = Table(data, repeatRows=1)

            # Style the table
            pdf_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.gray),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('INNERGRID', (0, 0), (-1, -1), 0.25, colors.black),
                ('BOX', (0, 0), (-1, -1), 0.25, colors.black),
            ]))

            # Measure table height
            table_height = pdf_table.wrap(usable_width, 0)[1]

            # Check if adding this table will exceed the page height
            if y_position - table_height < bottom_margin:
                c.showPage()  # Start a new page
                c.setFont("Helvetica", 12)
                y_position = pagesize[1] - top_margin

            # Draw the table
            pdf_table.drawOn(c, left_margin, y_position - table_height)

            # Adjust Y position after drawing table
            y_position -= table_height + 20  # Adding some space between tables

        # Save the canvas content into the BytesIO buffer
        c.showPage()
        c.save()

        # Write the PDF buffer to a file
        with open(pdf_path, 'wb') as f:
            f.write(buffer.getvalue())

        print(f"PDF created: {pdf_path}")
        return pdf_path

    except Exception as e:
        print(f"Error during PDF creation: {e}")
        return None
